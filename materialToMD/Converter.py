"""Conversion utilities to normalize many office formats into Markdown plus assets."""
from __future__ import annotations

import base64
import csv
import mimetypes
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
import re
import statistics
from typing import Iterable, List, Optional

import fitz  # PyMuPDF
from markdownify import markdownify as html_to_markdown
from openpyxl import load_workbook
from docx import Document
from docx.oxml.ns import qn
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class Asset:
    path: str  # relative path referenced in the Markdown output
    content_type: str
    data_base64: str  # base64-encoded asset bytes


@dataclass
class ConvertedDocument:
    markdown: str
    assets: List[Asset]


class ConversionError(Exception):
    """Raised when a file cannot be converted."""


_MD_ESCAPE_CHARS = {"|": "\\|", "\\": "\\\\", "`": "\\`"}
_ASSET_PREFIX = "assets"


def convert_to_markdown(file_bytes: bytes, filename: str) -> ConvertedDocument:
    """Convert a supported file's bytes into Markdown plus embedded assets."""
    extension = Path(filename).suffix.lower()
    assets: List[Asset] = []

    def add_asset(name: str, content: bytes, content_type: Optional[str] = None) -> str:
        asset_path = f"{_ASSET_PREFIX}/{name}"
        ctype = content_type or mimetypes.guess_type(name)[0] or "application/octet-stream"
        data_b64 = base64.b64encode(content).decode("ascii")
        assets.append(Asset(path=asset_path, content_type=ctype, data_base64=data_b64))
        return asset_path

    if extension == ".pdf":
        markdown = _convert_pdf(file_bytes, add_asset)
    elif extension in {".ppt", ".pptx"}:
        markdown = _convert_ppt(file_bytes, add_asset)
    elif extension == ".docx":
        markdown = _convert_docx(file_bytes, add_asset)
    elif extension == ".csv":
        markdown = _convert_csv(file_bytes)
    elif extension == ".xlsx":
        markdown = _convert_xlsx(file_bytes)
    elif extension == ".html":
        markdown = _convert_html(file_bytes)
    elif extension == ".txt":
        markdown = _convert_txt(file_bytes)
    else:
        raise ConversionError(f"Unsupported file type: {extension}")

    return ConvertedDocument(markdown=markdown.strip() + "\n", assets=assets)


def _convert_pdf(file_bytes: bytes, add_asset) -> str:
    """Extract text, headings, lists, code, and images from a PDF into Markdown."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    md_parts: List[str] = []
    image_index = 1

    for page_number in range(doc.page_count):
        page = doc.load_page(page_number)
        page_dict = page.get_text("dict")
        avg_font = _average_font_size(page_dict["blocks"])
        code_buffer: List[str] = []
        last_plain_index: Optional[int] = None

        def flush_code():
            nonlocal last_plain_index
            if code_buffer:
                md_parts.append("```\n" + "\n".join(code_buffer) + "\n```")
                code_buffer.clear()
                last_plain_index = None

        for block in page_dict["blocks"]:
            if block.get("type") == 0:
                rendered, is_code = _render_pdf_text_block(block, avg_font)
                if not rendered:
                    continue
                if isinstance(rendered, dict) and rendered.get("type") == "page_number":
                    page_number_text = rendered["text"]
                    try:
                        number_value = int(page_number_text)
                    except ValueError:
                        continue
                    if number_value == page_number + 1:
                        continue
                    if md_parts and md_parts[-1] == str(page_number_text):
                        continue
                    rendered = str(page_number_text)
                if is_code:
                    code_buffer.append(rendered)
                else:
                    flush_code()
                    if _is_plain_text(rendered):
                        if last_plain_index is not None:
                            md_parts[last_plain_index] += "\n" + rendered
                        else:
                            md_parts.append(rendered)
                            last_plain_index = len(md_parts) - 1
                    else:
                        last_plain_index = None
                        md_parts.append(rendered)
            elif block.get("type") == 1:
                flush_code()
                last_plain_index = None
                raw_image = block.get("image")
                if raw_image is None:
                    continue
                ext = block.get("ext") or "png"
                content_type = mimetypes.guess_type(f"file.{ext}")[0] or "image/png"
                if isinstance(raw_image, bytes):
                    image_bytes = raw_image
                    filename = f"page{page_number + 1}_img{image_index}.{ext}"
                else:
                    try:
                        extracted = doc.extract_image(raw_image)
                        image_bytes = extracted.get("image")
                        ext = extracted.get("ext", ext) or ext
                        content_type = mimetypes.guess_type(f"file.{ext}")[0] or "image/png"
                        filename = f"page{page_number + 1}_img{image_index}.{ext}"
                    except Exception:
                        pix = fitz.Pixmap(doc, raw_image)
                        # Normalize colorspace (including CMYK) and flatten alpha to avoid black fills.
                        if pix.colorspace and pix.colorspace.n not in (0, 1, 3):
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        if pix.alpha:
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        image_bytes = pix.tobytes("png")
                        ext = "png"
                        content_type = "image/png"
                        filename = f"page{page_number + 1}_img{image_index}.{ext}"
                asset_ref = add_asset(filename, image_bytes, content_type)
                md_parts.append(
                    f"![Page {page_number + 1} Image {image_index}]({asset_ref})"
                )
                image_index += 1

        flush_code()
        last_plain_index = None
        if md_parts and md_parts[-1] == "---":
            continue
        if page_number < doc.page_count - 1:
            md_parts.append("---")

    doc.close()
    markdown = "\n\n".join(part for part in md_parts if part).strip()
    return _collapse_blank_lines(markdown)


def _convert_docx(file_bytes: bytes, add_asset) -> str:
    """Render a DOCX document into Markdown including headings, lists, tables, and images."""
    document = Document(BytesIO(file_bytes))
    md_lines: List[str] = []
    image_index = 1

    def save_image(r_id: str) -> Optional[str]:
        nonlocal image_index
        if r_id not in document.part.related_parts:
            return None
        image_part = document.part.related_parts[r_id]
        ext = Path(str(image_part.partname)).suffix or ".png"
        filename = f"doc_image_{image_index}{ext}"
        image_index += 1
        return add_asset(filename, image_part.blob)

    def render_runs(runs) -> str:
        fragments: List[str] = []
        current_style: Optional[tuple[bool, bool, bool]] = None
        buffer = ""

        def flush_buffer():
            nonlocal buffer, current_style
            if not buffer:
                return
            bold, italic, underline = current_style or (False, False, False)
            wrapped = buffer
            if bold and italic:
                wrapped = f"***{wrapped}***"
            elif bold:
                wrapped = f"**{wrapped}**"
            elif italic:
                wrapped = f"*{wrapped}*"
            if underline:
                wrapped = f"__{wrapped}__"
            fragments.append(wrapped)
            buffer = ""

        for run in runs:
            blips = run.element.xpath(".//a:blip")
            if blips:
                flush_buffer()
                for blip in blips:
                    r_id = blip.get(qn("r:embed"))
                    rel = save_image(r_id) if r_id else None
                    if rel:
                        fragments.append(f"![Embedded Image]({rel})")
                current_style = None

            text = _escape_markdown(run.text or "")
            if not text:
                continue
            style = (bool(run.bold), bool(run.italic), bool(run.underline))
            if current_style != style:
                flush_buffer()
                current_style = style
            buffer += text

        flush_buffer()
        return "".join(fragments)

    def render_paragraph(paragraph: Paragraph) -> str:
        style = paragraph.style.name if paragraph.style else ""
        content = render_runs(paragraph.runs).strip()
        if not content:
            return ""

        if style.startswith("Heading"):
            level = _heading_level(style)
            return _format_heading(content, level)

        list_level = _list_level(paragraph)
        if list_level is not None:
            indent = "  " * list_level
            return f"{indent}- {_wrap_direction(content, inline=True)}"

        if "List" in style:
            return f"- {_wrap_direction(content, inline=True)}"

        if style.lower().startswith("code"):
            return f"```\n{content}\n```"

        return _wrap_direction(content, inline=False)

    for block in _iter_block_items(document):
        if isinstance(block, Paragraph):
            rendered = render_paragraph(block)
            if rendered:
                md_lines.append(rendered)
                md_lines.append("")
        elif isinstance(block, Table):
            md_lines.append(_table_to_markdown(block, render_runs))
            md_lines.append("")

    return "\n".join(md_lines).strip()


def _iter_block_items(parent) -> Iterable[Paragraph | Table]:
    """Yield block-level elements (paragraphs and tables) from a DOCX parent."""
    body = parent.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _list_level(paragraph: Paragraph) -> Optional[int]:
    """Determine the nesting level for numbered/bulleted paragraphs in DOCX."""
    p = paragraph._p
    if p is None or p.pPr is None or p.pPr.numPr is None:
        return None
    ilvl = p.pPr.numPr.ilvl
    return int(ilvl.val) if ilvl is not None else 0


def _table_to_markdown(table: Table, render_runs) -> str:
    """Convert a DOCX table to Markdown syntax."""
    rows: List[List[str]] = []
    for row in table.rows:
        cells: List[str] = []
        for cell in row.cells:
            cell_fragments = [
                render_runs(paragraph.runs).strip() for paragraph in cell.paragraphs
            ]
            cell_text = "\n".join(filter(None, cell_fragments)) or " "
            cell_text = _wrap_direction(cell_text.replace("\n", "<br>"), inline=True)
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    header_line = "| " + " | ".join(_escape_markdown(cell) for cell in header) + " |"
    divider = "| " + " | ".join("---" for _ in header) + " |"
    body_lines = [
        "| " + " | ".join(_escape_markdown(cell) for cell in row) + " |"
        for row in body
    ]
    return "\n".join([header_line, divider] + body_lines)


def _convert_ppt(file_bytes: bytes, add_asset) -> str:
    """Render PPT/PPTX slides into Markdown with headings, bullets, and images."""
    pres = Presentation(BytesIO(file_bytes))
    md_parts: List[str] = []
    image_index = 1

    def save_image(image, slide_number: int) -> Optional[str]:
        nonlocal image_index
        ext = image.ext or ".png"
        filename = f"slide{slide_number}_img{image_index}.{ext.lstrip('.')}"
        image_index += 1
        content_type = image.content_type if hasattr(image, "content_type") else None
        return add_asset(filename, image.blob, content_type)

    for slide_number, slide in enumerate(pres.slides, start=1):
        md_parts.append(_format_heading(f"Slide {slide_number}", 1))
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                rel = save_image(shape.image, slide_number)
                if rel:
                    md_parts.append(f"![Slide {slide_number} Image]({rel})")
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = _escape_markdown(paragraph.text or "").strip()
                if not text:
                    continue
                prefix = "- " * (paragraph.level + 1) if paragraph.level >= 0 else ""
                md_parts.append(f"{prefix}{_wrap_direction(text, inline=True)}")
        md_parts.append("")

    return "\n".join(md_parts).strip()


def _convert_csv(file_bytes: bytes) -> str:
    """Convert CSV content into a Markdown table."""
    text = file_bytes.decode("utf-8", errors="ignore")
    reader = csv.reader(text.splitlines())
    rows = list(reader)
    if not rows:
        return ""
    header = rows[0]
    width = len(header)
    md_lines = [
        "| " + " | ".join(_escape_markdown(cell.strip()) for cell in header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        padded = list(row) + [""] * (width - len(row))
        md_lines.append("| " + " | ".join(_escape_markdown(cell.strip() or " ") for cell in padded) + " |")
    return "\n".join(md_lines)


def _convert_xlsx(file_bytes: bytes) -> str:
    """Convert each worksheet of an Excel workbook into Markdown tables."""
    wb = load_workbook(filename=BytesIO(file_bytes), data_only=True)
    parts: List[str] = []
    for sheet in wb.worksheets:
        title = sheet.title or "Sheet"
        parts.append(_format_heading(title, 2))
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        header = ["" if cell is None else str(cell) for cell in rows[0]]
        width = len(header)
        table_lines = [
            "| " + " | ".join(_escape_markdown(str(cell)) for cell in header) + " |",
            "| " + " | ".join("---" for _ in header) + " |",
        ]
        for row in rows[1:]:
            cells = list(row) if row else []
            padded = cells + ["" for _ in range(width - len(cells))]
            table_lines.append(
                "| " + " | ".join(_escape_markdown("" if cell is None else str(cell)) for cell in padded) + " |"
            )
        parts.append("\n".join(table_lines))
    return "\n\n".join(part for part in parts if part).strip()


def _convert_html(file_bytes: bytes) -> str:
    """Convert HTML content to Markdown using markdownify."""
    html = file_bytes.decode("utf-8", errors="ignore")
    markdown = html_to_markdown(html, heading_style="ATX")
    return _collapse_blank_lines(markdown)


def _convert_txt(file_bytes: bytes) -> str:
    """Normalize plain text by collapsing whitespace."""
    text = file_bytes.decode("utf-8", errors="ignore")
    return _collapse_blank_lines(text)


def _heading_level(style_name: str) -> int:
    """Infer heading level from a DOCX style name."""
    for token in style_name.split():
        if token.isdigit():
            return int(token)
        if token and token[-1].isdigit():
            try:
                return int(token[-1])
            except ValueError:
                continue
    return 1


def _escape_markdown(text: str) -> str:
    """Escape special Markdown characters to prevent unintended formatting."""
    escaped = text
    for char, replacement in _MD_ESCAPE_CHARS.items():
        escaped = escaped.replace(char, replacement)
    return escaped


def _average_font_size(blocks) -> float:
    """Compute an average font size for a PDF page to help detect headings."""
    sizes: List[float] = []
    for block in blocks:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                size = span.get("size")
                if size:
                    sizes.append(size)
    return statistics.mean(sizes) if sizes else 12.0


def _render_pdf_text_block(block, avg_font: float) -> tuple[str | dict, bool]:
    """Render a PDF text block to Markdown and flag whether it looks like code."""
    lines: List[str] = []
    font_sizes: List[float] = []
    for line in block.get("lines", []):
        fragments = [span.get("text", "") for span in line.get("spans", [])]
        text = " ".join(fragment for fragment in fragments if fragment)
        text = _normalize_whitespace(text)
        if not text:
            continue
        lines.append(text)
        for span in line.get("spans", []):
            size = span.get("size")
            if size:
                font_sizes.append(size)

    lines = _merge_numbered_lines(lines)

    if not lines:
        return "", False

    if _looks_like_code_block(lines):
        code = "\n".join(lines)
        return code, True

    if all(_looks_like_list_item(line) for line in lines):
        return "\n".join(_format_list_line(line) for line in lines), False

    if _looks_like_page_number(lines):
        return {"type": "page_number", "text": lines[0]}, False

    block_font = max(font_sizes) if font_sizes else avg_font
    text = _normalize_whitespace(" ".join(lines))
    heading_level = _pdf_heading_level(block_font, avg_font)
    if heading_level:
        return _format_heading(text, heading_level), False
    return _wrap_direction(text, inline=False), False


_BULLET_PATTERN = re.compile(r"^\s*([\-\u2022\u2023\u25CF\*]|(\d+[\.\)]))\s+")


def _looks_like_list_item(line: str) -> bool:
    """Check whether a line begins with a bullet or numbered prefix."""
    return bool(_BULLET_PATTERN.match(line))


def _format_list_line(line: str) -> str:
    """Normalize a list line to Markdown bullet syntax."""
    stripped = _BULLET_PATTERN.sub("", line).strip()
    return f"- {_wrap_direction(stripped, inline=True)}"


def _normalize_whitespace(text: str) -> str:
    """Collapse whitespace and insert bidi spacing where necessary."""
    collapsed = re.sub(r"\s+", " ", text or "").strip()
    return _insert_bidi_spacing(collapsed)


def _pdf_heading_level(block_font: float, avg_font: float) -> int:
    """Determine if a PDF text block should be treated as a heading based on font size."""
    if block_font >= avg_font * 1.8:
        return 1
    if block_font >= avg_font * 1.35:
        return 2
    return 0


_CODE_HINT = re.compile(r"[{}<>;=;()]")


def _looks_like_code_block(lines: List[str]) -> bool:
    """Detect if a collection of lines resembles a code block."""
    if not lines:
        return False
    code_like = sum(1 for line in lines if _CODE_HINT.search(line))
    if len(lines) == 1:
        return code_like == 1
    return code_like >= max(2, len(lines) // 2)


_HEBREW_PATTERN = re.compile(r"[\u0590-\u05FF]")
_LTR_CHARS = "A-Za-z0-9"
_BIDI_BOUNDARY = re.compile(
    rf"(?<=[\u0590-\u05FF])(?=[{_LTR_CHARS}])|(?<=[{_LTR_CHARS}])(?=[\u0590-\u05FF])"
)
_LINE_NUMBER_PATTERN = re.compile(r"^\d+\.?$")


def _is_rtl(text: str) -> bool:
    """Return True if the text contains right-to-left characters (Hebrew)."""
    return bool(_HEBREW_PATTERN.search(text or ""))


def _wrap_direction(text: str, inline: bool) -> str:
    """Wrap RTL text in direction-aware tags when needed."""
    if not text or not _is_rtl(text):
        return text
    tag = "span" if inline else "div"
    return f"<{tag} dir=\"rtl\">{text}</{tag}>"


def _insert_bidi_spacing(text: str) -> str:
    """Insert a space between RTL and LTR text boundaries for readability."""
    return _BIDI_BOUNDARY.sub(" ", text)


def _collapse_blank_lines(text: str) -> str:
    """Collapse multiple consecutive blank lines to a single blank line."""
    lines = text.splitlines()
    collapsed: List[str] = []
    previous_blank = False
    for line in lines:
        if line.strip():
            collapsed.append(line)
            previous_blank = False
        else:
            if not previous_blank:
                collapsed.append("")
                previous_blank = True
    return "\n".join(collapsed).strip()


def _is_plain_text(text: str) -> bool:
    """Heuristic to decide if Markdown text is plain prose for concatenation."""
    if not text:
        return False
    stripped = text.lstrip()
    return not (
        stripped.startswith(("<", "#", "-", "!", "`", "|", "*"))
        or stripped.lower().startswith("```")
    )


def _looks_like_page_number(lines: List[str]) -> bool:
    """Detect isolated page number strings to avoid rendering them."""
    if len(lines) != 1:
        return False
    text = lines[0].strip()
    return text.isdigit() and len(text) <= 2


def _merge_numbered_lines(lines: List[str]) -> List[str]:
    """Merge line-number-only rows with the following line for better flow."""
    merged: List[str] = []
    i = 0
    while i < len(lines):
        line = lines[i]
        if _LINE_NUMBER_PATTERN.match(line.strip()) and i + 1 < len(lines):
            merged.append(f"{line.strip()} {lines[i + 1].strip()}")
            i += 2
        else:
            merged.append(line)
            i += 1
    return merged


def _format_heading(text: str, level: int) -> str:
    """Format a string as a Markdown heading respecting RTL content."""
    text = text.strip()
    if not text:
        return ""
    if _is_rtl(text):
        return f"<h{level} dir=\"rtl\">{text}</h{level}>"
    return f"{'#' * level} {text}"


__all__ = ["convert_to_markdown", "ConversionError", "ConvertedDocument", "Asset"]
