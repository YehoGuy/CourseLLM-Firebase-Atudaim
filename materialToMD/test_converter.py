import base64
from io import BytesIO
from pathlib import Path

import fitz
import pytest
from PIL import Image, ImageDraw
from docx import Document
from docx.shared import Inches
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches as PptInches

from Converter import ConversionError, convert_to_markdown
from app import app


def _make_sample_image(tmp_path: Path) -> Path:
    path = tmp_path / "sample.png"
    img = Image.new("RGB", (120, 80), color=(20, 30, 200))
    draw = ImageDraw.Draw(img)
    draw.text((10, 30), "IMG", fill=(255, 255, 255))
    img.save(path)
    return path


def test_docx_to_markdown(tmp_path: Path) -> None:
    image_path = _make_sample_image(tmp_path)

    doc = Document()
    doc.add_heading("Docx Title", level=1)
    doc.add_heading("שלום כיתה א", level=2)
    para = doc.add_paragraph("Intro text with ")
    run = para.add_run("bold italic")
    run.bold = True
    run.italic = True
    doc.add_paragraph("Bullet entry", style="List Bullet")
    doc.add_paragraph("שלום עולם")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "C1"
    table.cell(1, 1).text = "C2"

    doc.add_picture(str(image_path), width=Inches(1.0))
    buff = BytesIO()
    doc.save(buff)

    result = convert_to_markdown(buff.getvalue(), "sample.docx")
    md_text = result.markdown

    assert "# Docx Title" in md_text
    assert "***bold italic***" in md_text
    assert "- Bullet entry" in md_text
    assert "| H1 | H2 |" in md_text
    assert "Embedded Image" in md_text
    assert "<h2 dir=\"rtl\">שלום כיתה א</h2>" in md_text
    assert "<div dir=\"rtl\">שלום עולם</div>" in md_text
    assert result.assets
    img_asset = next(a for a in result.assets if "doc_image" in a.path)
    img_bytes = base64.b64decode(img_asset.data_base64)
    assert img_bytes.startswith(b"\x89PNG")


def test_ppt_to_markdown(tmp_path: Path) -> None:
    image_path = _make_sample_image(tmp_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    left = top = PptInches(1)
    width = height = PptInches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = "Slide Title"
    para = textbox.text_frame.add_paragraph()
    para.text = "Indented bullet"
    para.level = 1
    slide.shapes.add_picture(str(image_path), left, PptInches(4), width=PptInches(2), height=PptInches(1.5))
    buff = BytesIO()
    prs.save(buff)

    result = convert_to_markdown(buff.getvalue(), "sample.pptx")
    md_text = result.markdown

    assert "# Slide 1" in md_text
    assert "Slide Title" in md_text
    assert "- - Indented bullet" in md_text  # two levels of bullets
    assert "![Slide 1 Image]" in md_text
    assert result.assets


def test_pdf_to_markdown(tmp_path: Path) -> None:
    image_path = _make_sample_image(tmp_path)

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "PDF Heading")
    page.insert_text((72, 100), "Second line for PDF")
    page.insert_text((72, 130), "שלום area")
    rect = fitz.Rect(72, 140, 172, 220)
    page.insert_image(rect, filename=str(image_path))
    code_lines = [
        "public class Example {",
        "  private int value;",
        "  public void set(int v) {",
        "    this.value = v;",
        "  }",
        "}",
    ]
    y = 240
    for line in code_lines:
        page.insert_text((72, y), line)
        y += 18
    pdf_bytes = doc.tobytes()
    doc.close()

    result = convert_to_markdown(pdf_bytes, "sample.pdf")
    md_text = result.markdown

    assert "PDF Heading" in md_text
    assert "Second line for PDF" in md_text
    assert "Image 1" in md_text
    assert result.assets
    assert md_text.count("```") == 2  # one pair of fences


def test_unsupported_extension(tmp_path: Path) -> None:
    src = tmp_path / "note.txt"
    src.write_text("hello", encoding="utf-8")
    with pytest.raises(ConversionError):
        convert_to_markdown(src.read_bytes(), src.name)


def test_fastapi_convert_endpoint(tmp_path: Path) -> None:
    client = TestClient(app)
    image_path = _make_sample_image(tmp_path)

    doc = Document()
    doc.add_paragraph("API Docx")
    doc.add_picture(str(image_path))
    buff = BytesIO()
    doc.save(buff)

    files = {"file": ("api.docx", buff.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
    resp = client.post("/convert", files=files)

    assert resp.status_code == 200
    payload = resp.json()
    assert "markdown" in payload and "assets" in payload
    assert "API Docx" in payload["markdown"]
    assert payload["assets"]
    first_asset = payload["assets"][0]
    decoded = base64.b64decode(first_asset["data_base64"])
    assert decoded.startswith(b"\x89PNG") or decoded.startswith(b"\xff\xd8")


def test_fastapi_rejects_unknown_extension(tmp_path: Path) -> None:
    client = TestClient(app)
    files = {"file": ("note.txt", b"hello", "text/plain")}
    resp = client.post("/convert", files=files)
    assert resp.status_code == 400
