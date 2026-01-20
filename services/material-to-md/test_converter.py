import base64
from io import BytesIO
from pathlib import Path

import fitz
import pytest
from PIL import Image, ImageDraw
from docx import Document
from docx.shared import Inches
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PptInches

from Converter import Asset, ConversionError, convert_to_markdown
from app import create_app
from settings import Settings
from storage import LocalStorageBackend


def _test_settings(tmp_path: Path) -> Settings:
    return Settings(
        db_url=f"sqlite:///{tmp_path}/test.db",
        storage_root=tmp_path / "storage",
        max_concurrent_jobs=1,
        enable_scheduler=False,
        retry_initial_delay=0.01,
        retry_backoff_factor=1.0,
        max_retries=1,
        polling_interval=1,
    )


def _make_sample_image(tmp_path: Path) -> Path:
    path = tmp_path / "sample.png"
    img = Image.new("RGB", (120, 80), color=(20, 30, 200))
    draw = ImageDraw.Draw(img)
    draw.text((10, 30), "IMG", fill=(255, 255, 255))
    img.save(path)
    return path


def test_docx_to_markdown(tmp_path: Path) -> None:
    image_path = _make_sample_image(tmp_path)

    doc = Document()
    doc.add_heading("Docx Title", level=1)
    doc.add_heading("שלום כיתה א", level=2)
    para = doc.add_paragraph("Intro text with ")
    run = para.add_run("bold italic")
    run.bold = True
    run.italic = True
    doc.add_paragraph("Bullet entry", style="List Bullet")
    doc.add_paragraph("שלום עולם")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "C1"
    table.cell(1, 1).text = "C2"

    doc.add_picture(str(image_path), width=Inches(1.0))
    buff = BytesIO()
    doc.save(buff)

    result = convert_to_markdown(buff.getvalue(), "sample.docx")
    md_text = result.markdown

    assert "# Docx Title" in md_text
    assert "***bold italic***" in md_text
    assert "- Bullet entry" in md_text
    assert "| H1 | H2 |" in md_text
    assert "Embedded Image" in md_text
    assert "<h2 dir=\"rtl\">שלום כיתה א</h2>" in md_text
    assert "<div dir=\"rtl\">שלום עולם</div>" in md_text
    assert result.assets
    img_asset = next(a for a in result.assets if "doc_image" in a.path)
    img_bytes = base64.b64decode(img_asset.data_base64)
    assert img_bytes.startswith(b"\x89PNG")


def test_ppt_to_markdown(tmp_path: Path) -> None:
    image_path = _make_sample_image(tmp_path)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    left = top = PptInches(1)
    width = height = PptInches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = "Slide Title"
    para = textbox.text_frame.add_paragraph()
    para.text = "Indented bullet"
    para.level = 1
    slide.shapes.add_picture(str(image_path), left, PptInches(4), width=PptInches(2), height=PptInches(1.5))
    buff = BytesIO()
    prs.save(buff)

    result = convert_to_markdown(buff.getvalue(), "sample.pptx")
    md_text = result.markdown

    assert "# Slide 1" in md_text
    assert "Slide Title" in md_text
    assert "- - Indented bullet" in md_text  # two levels of bullets
    assert "![Slide 1 Image]" in md_text
    assert result.assets


def test_pdf_to_markdown(tmp_path: Path) -> None:
    image_path = _make_sample_image(tmp_path)

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "PDF Heading")
    page.insert_text((72, 100), "Second line for PDF")
    page.insert_text((72, 130), "שלום area")
    rect = fitz.Rect(72, 140, 172, 220)
    page.insert_image(rect, filename=str(image_path))
    code_lines = [
        "public class Example {",
        "  private int value;",
        "  public void set(int v) {",
        "    this.value = v;",
        "  }",
        "}",
    ]
    y = 240
    for line in code_lines:
        page.insert_text((72, y), line)
        y += 18
    pdf_bytes = doc.tobytes()
    doc.close()

    result = convert_to_markdown(pdf_bytes, "sample.pdf")
    md_text = result.markdown

    assert "PDF Heading" in md_text
    assert "Second line for PDF" in md_text
    assert "Image 1" in md_text
    assert result.assets
    assert md_text.count("```") == 2  # one pair of fences


def test_csv_to_markdown(tmp_path: Path) -> None:
    csv_path = tmp_path / "sample.csv"
    csv_path.write_text("name,age\nAlice,30\nBob,25", encoding="utf-8")
    result = convert_to_markdown(csv_path.read_bytes(), csv_path.name)
    md = result.markdown
    assert "| name | age |" in md
    assert "| Alice | 30 |" in md


def test_xlsx_to_markdown(tmp_path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet One"
    ws.append(["col1", "col2"])
    ws.append(["val1", "val2"])
    buff = BytesIO()
    wb.save(buff)

    result = convert_to_markdown(buff.getvalue(), "workbook.xlsx")
    md = result.markdown
    assert "Sheet One" in md
    assert "| col1 | col2 |" in md
    assert "| val1 | val2 |" in md


def test_html_and_txt_conversion(tmp_path: Path) -> None:
    html = "<html><body><h1>Title</h1><p>Paragraph</p></body></html>"
    html_result = convert_to_markdown(html.encode("utf-8"), "page.html")
    assert "# Title" in html_result.markdown
    assert "Paragraph" in html_result.markdown

    txt_path = tmp_path / "sample.txt"
    txt_path.write_text("first line\n\nsecond line", encoding="utf-8")
    txt_result = convert_to_markdown(txt_path.read_bytes(), txt_path.name)
    assert "first line" in txt_result.markdown
    assert "second line" in txt_result.markdown


def test_txt_collapses_blank_lines(tmp_path: Path) -> None:
    txt_path = tmp_path / "multi.txt"
    txt_path.write_text("line1\n\n\n\nline2\n\n\nline3", encoding="utf-8")
    txt_result = convert_to_markdown(txt_path.read_bytes(), txt_path.name)
    assert "\n\n\n" not in txt_result.markdown
    assert "line1" in txt_result.markdown
    assert "line2" in txt_result.markdown
    assert "line3" in txt_result.markdown


def test_csv_escapes_markdown_chars(tmp_path: Path) -> None:
    csv_path = tmp_path / "escape.csv"
    csv_path.write_text("col1,col2\npipe|here,back`tick\\\\test", encoding="utf-8")
    result = convert_to_markdown(csv_path.read_bytes(), csv_path.name)
    md = result.markdown
    assert "pipe\\\\|here" in md
    assert "back\\`tick\\\\\\\\test" in md


def test_xlsx_includes_empty_sheet_heading(tmp_path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Has Data"
    ws.append(["h1", "h2"])
    ws.append(["v1", "v2"])
    wb.create_sheet("Empty Sheet")
    buff = BytesIO()
    wb.save(buff)

    result = convert_to_markdown(buff.getvalue(), "multi.xlsx")
    md = result.markdown
    assert "## Has Data" in md
    assert "| h1 | h2 |" in md
    assert "## Empty Sheet" in md


def test_unsupported_extension(tmp_path: Path) -> None:
    src = tmp_path / "note.bin"
    src.write_bytes(b"hello")
    with pytest.raises(ConversionError):
        convert_to_markdown(src.read_bytes(), src.name)


def test_fastapi_convert_endpoint(tmp_path: Path) -> None:
    client_app = create_app(settings=_test_settings(tmp_path))
    with TestClient(client_app) as client:
        image_path = _make_sample_image(tmp_path)

        doc = Document()
        doc.add_paragraph("API Docx")
        doc.add_picture(str(image_path))
        buff = BytesIO()
        doc.save(buff)

        files = {
            "file": (
                "api.docx",
                buff.getvalue(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        }
        resp = client.post("/convert", files=files)

        assert resp.status_code == 200
        payload = resp.json()
        assert "markdown" in payload and "assets" in payload
        assert "API Docx" in payload["markdown"]
        assert payload["assets"]
        first_asset = payload["assets"][0]
        decoded = base64.b64decode(first_asset["data_base64"])
        assert decoded.startswith(b"\x89PNG") or decoded.startswith(b"\xff\xd8")


def test_fastapi_rejects_unknown_extension(tmp_path: Path) -> None:
    client_app = create_app(settings=_test_settings(tmp_path))
    with TestClient(client_app) as client:
        files = {"file": ("note.bin", b"hello", "application/octet-stream")}
        resp = client.post("/convert", files=files)
        assert resp.status_code == 400


def test_storage_writes_assets(tmp_path: Path) -> None:
    settings = _test_settings(tmp_path)
    storage = LocalStorageBackend(settings)
    storage.ensure_dirs()
    encoded = base64.b64encode(b"img-bytes").decode("ascii")
    asset = Asset(path="assets/test.png", content_type="image/png", data_base64=encoded)

    output_path = storage.write_processed("note.txt", "# Hello\n", [asset])
    assert Path(output_path).exists()
    processed_dir = Path(settings.storage_root) / settings.processed_dir
    assert (processed_dir / "note.md").exists()
    assert (processed_dir / "assets" / "test.png").exists()
